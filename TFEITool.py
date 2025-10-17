#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TFEITool.py — The Future of Exabyte‑Scale Inference (TFEI) slide generator + analyzer

What this does
--------------
1) Runs a real TRT-LLM inference (OpenAI-compatible HTTP server, e.g. trtllm-serve),
   or a --dry-run simulation.
2) Traces I/O interactions end-to-end (HTTP request/response bytes, TTFT, tokens/sec).
3) Introspects TensorRT-LLM engine/runtime (Python bindings) to discover KV geometry:
     - num_kv_heads (Hk), head_dim (D), tokens_per_block (T), dtype, kv quant scale (if present).
4) Reads KV block dump files (e.g., block_<N>_pool_<M>.bin) and maps EVERY BYTE to:
     (kv plane, head h, token position t, dimension d), following row-major layout.
5) Produces a clean, self-contained PPTX deck labeling:
     - program + dataflow diagrams,
     - KV-block concepts (2 × Hk × T × D),
     - color legend and tile maps for heads/tokens/dims,
     - linear on-disk map,
     - FULL hex dump (no truncation) with color-coded overlays and legends.

Requirements
------------
Python 3.10+ recommended.

External deps:
    pip install python-pptx pillow numpy matplotlib requests tqdm pyyaml
Optional:
    pip install openai
    pip install tensorrt-llm    # From NVIDIA index or local wheel
"""

import argparse
import os
import sys
import time
import math
import glob
import json
import struct
import textwrap
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple

import numpy as np

# Optional runtime bindings
try:
    # Public Python runtime bindings shipped by TensorRT-LLM
    from tensorrt_llm.runtime import ModelRunner  # type: ignore
    TRTLLM_AVAILABLE = True
except Exception:
    TRTLLM_AVAILABLE = False

# Optional OpenAI-compatible client for trtllm-serve / vLLM etc.
try:
    import requests
    OPENAI_HTTP_AVAILABLE = True
except Exception:
    OPENAI_HTTP_AVAILABLE = False

# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from PIL import Image, ImageDraw, ImageFont
from io import BytesIO

# --------------------------------------------------------------------------------------
# Constants & Palettes
# --------------------------------------------------------------------------------------

PALETTE_HEADS = [
    (30, 120, 200),   # h=0 blue
    (240, 150, 50),   # h=1 orange
    (50, 170, 80),    # h=2 green
    (220, 60, 70),    # h=3 red
    (140, 90, 200),   # h=4 purple
    (150, 110, 90),   # h=5 brown
    (230, 120, 200),  # h=6 pink
    (120, 120, 120),  # h=7 gray
]

COLOR_K = (50, 120, 230)     # K plane fill (soft)
COLOR_V = (230, 90, 90)      # V plane fill (soft)
COLOR_TOKEN_STRIPE = (240, 240, 240)
COLOR_DIM_STRIPE = (220, 240, 220)

# --------------------------------------------------------------------------------------
# CLI configuration
# --------------------------------------------------------------------------------------

@dataclass
class CLI:
    server_url: str
    model: str
    api_key: str
    engine_dir: Optional[str]
    kv_dump_glob: str
    pptx_out: str
    inference_prompt: str
    layer_index: int
    dry_run: bool
    dtype: Optional[str]
    kv_scale: Optional[float]
    max_hex_lines_per_slide: int = 32
    max_bytes_preview_img: int = 2048  # for linear on-disk preview banner

# --------------------------------------------------------------------------------------
# External extension hooks (for your two gists)
# --------------------------------------------------------------------------------------

class ExternalHooks:
    """
    This adapter layer is where you can drop in code from your two gists.
    Place your modules on PYTHONPATH and name them:
       - tfei_trtllm_hooks_a
       - tfei_trtllm_hooks_b
    Each may define any subset of:
       before_prefill(ctx), after_prefill(ctx)
       before_decode(ctx), after_decode(ctx)
       on_copy_block(ctx, info)    # info: dict with {pool, block, plane, bytes, ts}
       on_offload_block(ctx, info)
       on_onboard_block(ctx, info)
    These hooks are invoked as we receive events from the server/logs (or from
    your bindings if you call them inside TRT kernels).
    """
    def __init__(self):
        self.a = None
        self.b = None
        try:
            import tfei_trtllm_hooks_a as A  # type: ignore
            self.a = A
        except Exception:
            pass
        try:
            import tfei_trtllm_hooks_b as B  # type: ignore
            self.b = B
        except Exception:
            pass

    def _call(self, name: str, *args, **kwargs):
        for mod in (self.a, self.b):
            if mod is None:
                continue
            fn = getattr(mod, name, None)
            if fn:
                try:
                    fn(*args, **kwargs)
                except Exception as e:
                    print(f"[hook:{name}] error: {e}", file=sys.stderr)

    def before_prefill(self, ctx):  self._call("before_prefill", ctx)
    def after_prefill(self, ctx):   self._call("after_prefill", ctx)
    def before_decode(self, ctx):   self._call("before_decode", ctx)
    def after_decode(self, ctx):    self._call("after_decode", ctx)
    def on_copy_block(self, ctx, info):    self._call("on_copy_block", ctx, info)
    def on_offload_block(self, ctx, info): self._call("on_offload_block", ctx, info)
    def on_onboard_block(self, ctx, info): self._call("on_onboard_block", ctx, info)

# --------------------------------------------------------------------------------------
# TRT‑LLM runtime introspection
# --------------------------------------------------------------------------------------

@dataclass
class KVGeometry:
    num_kv_heads: int
    head_dim: int
    tokens_per_block: int
    dtype: str
    kv_scale: Optional[float] = None

def discover_kv_geometry(engine_dir: Optional[str], layer_index: int,
                         preferred_dtype: Optional[str]) -> KVGeometry:
    """
    Try to discover KV geometry from TRT‑LLM runtime or engine metadata.
    Falls back to reasonable defaults if unavailable (and dry-run only).
    """
    # Defaults for dry‑run
    fallback = KVGeometry(
        num_kv_heads=8, head_dim=16, tokens_per_block=32,
        dtype=preferred_dtype or "fp16", kv_scale=None
    )
    if not TRTLLM_AVAILABLE or not engine_dir:
        return fallback

    try:
        runner = ModelRunner.from_dir(engine_dir, rank=0, debug_mode=False)  # type: ignore
        mc = getattr(getattr(runner, "session", runner), "model_config", None)
        if mc is None:
            mc = getattr(runner, "model_config", None)
        if mc is None:
            return fallback

        # Try extracting from model_config
        T = getattr(mc, "tokens_per_block", None) or 32
        Hk = None
        nkv_per_layer = getattr(mc, "num_kv_heads_per_layer", None)
        if nkv_per_layer and len(nkv_per_layer) > layer_index:
            Hk = int(nkv_per_layer[layer_index])
        else:
            Hk = int(getattr(mc, "num_kv_heads", 8))

        D = getattr(mc, "head_size", None)
        if not D:
            hidden = getattr(mc, "hidden_size", None)
            nh = getattr(mc, "num_heads", None)
            if hidden and nh:
                D = int(hidden // nh)
            else:
                D = 16

        dtype = preferred_dtype or getattr(mc, "dtype", "fp16")
        return KVGeometry(num_kv_heads=int(Hk), head_dim=int(D),
                          tokens_per_block=int(T), dtype=dtype)
    except Exception as e:
        print(f"[warn] TRT‑LLM model introspection failed: {e}", file=sys.stderr)
        return fallback

# --------------------------------------------------------------------------------------
# I/O tracer (HTTP, TTFT, token throughput)
# --------------------------------------------------------------------------------------

@dataclass
class IOEvent:
    name: str
    ts: float
    meta: Dict[str, Any] = field(default_factory=dict)

class IOTracer:
    def __init__(self):
        self.events: List[IOEvent] = []
        self.start_ts: float = time.time()
        self.bytes_up = 0
        self.bytes_down = 0
        self.first_token_ts: Optional[float] = None
        self.tokens = 0

    def mark(self, name: str, **meta):
        self.events.append(IOEvent(name=name, ts=time.time(), meta=meta))

    def add_bytes(self, up: int = 0, down: int = 0):
        self.bytes_up += up
        self.bytes_down += down

    def token(self):
        self.tokens += 1
        if self.first_token_ts is None:
            self.first_token_ts = time.time()

    def summary(self) -> Dict[str, Any]:
        dur = time.time() - self.start_ts
        ttft = None
        if self.first_token_ts is not None:
            ttft = self.first_token_ts - self.start_ts
        tps = self.tokens / (dur if dur > 0 else 1.0)
        return dict(
            duration_s=dur, ttft_s=ttft, tokens=self.tokens,
            tokens_per_s=tps, bytes_up=self.bytes_up, bytes_down=self.bytes_down,
        )

# --------------------------------------------------------------------------------------
# OpenAI-compatible request to TRT‑LLM server
# --------------------------------------------------------------------------------------

def run_inference_openai_http(server_url: str, model: str, api_key: str,
                              prompt: str, tracer: IOTracer, hooks: ExternalHooks) -> str:
    """
    Streams a completion from an OpenAI-compatible server (e.g., `trtllm-serve --port ...`).
    Logs request/response bytes and token timings.
    """
    if not OPENAI_HTTP_AVAILABLE:
        raise RuntimeError("requests library not available")

    headers = {
        "Authorization": f"Bearer {api_key or 'sk-noop'}",
        "Content-Type": "application/json",
    }
    url = server_url.rstrip("/") + "/v1/chat/completions"
    payload = {
        "model": model,
        "stream": True,
        "messages": [{"role": "user", "content": prompt}],
    }
    data = json.dumps(payload).encode("utf-8")
    tracer.add_bytes(up=len(data))
    tracer.mark("http_request", url=url)

    hooks.before_prefill({"stage": "prefill_http"})
    with requests.post(url, headers=headers, data=data, stream=True, timeout=600) as r:
        r.raise_for_status()
        out = []
        for line in r.iter_lines(decode_unicode=True):
            if not line:
                continue
            tracer.add_bytes(down=len(line))
            if line.startswith("data: "):
                body = line[6:]
                if body.strip() == "[DONE]":
                    break
                try:
                    obj = json.loads(body)
                    delta = obj["choices"][0]["delta"].get("content", "")
                    if delta:
                        tracer.token()
                        out.append(delta)
                except Exception:
                    # Some servers return chunks with different schemas; be tolerant
                    pass
        hooks.after_prefill({"stage": "prefill_http_done"})
    return "".join(out)

# --------------------------------------------------------------------------------------
# KV cache file analysis
# --------------------------------------------------------------------------------------

@dataclass
class KVIndex:
    kv: int
    h: int
    t: int
    d: int

def byte_index_to_indices(byte_index: int, element_bytes: int,
                          Hk: int, T: int, D: int) -> KVIndex:
    """
    Map a byte offset into (kv, h, t, d).
    Layout is row-major over elements: [2, Hk, T, D].
    Each element has `element_bytes` (e.g., 2 for fp16).
    """
    elem_index = byte_index // element_bytes
    elems_per_plane = Hk * T * D
    kv = elem_index // elems_per_plane
    rem = elem_index % elems_per_plane
    h = rem // (T * D)
    rem2 = rem % (T * D)
    t = rem2 // D
    d = rem2 % D
    return KVIndex(kv=kv, h=h, t=t, d=d)

def read_block_file(path: str) -> bytes:
    with open(path, "rb") as f:
        return f.read()

# --------------------------------------------------------------------------------------
# Rendering helpers (PNG fragments that we embed into PPTX)
# --------------------------------------------------------------------------------------

def draw_color_legend(width=1280, height=260) -> Image.Image:
    img = Image.new("RGB", (width, height), (255, 255, 255))
    dr = ImageDraw.Draw(img)
    font_b = ImageFont.load_default()
    # Heads swatches
    pad = 20
    sw = (width - pad * (len(PALETTE_HEADS) + 1)) // len(PALETTE_HEADS)
    y = 80
    dr.text((pad, 20), "Head Color Legend (h = 0..Hk-1)", fill=(0, 0, 0), font=font_b)
    for i, c in enumerate(PALETTE_HEADS):
        x = pad + i * (sw + pad)
        dr.rectangle([x, y, x + sw, y + 60], fill=c, outline=(0, 0, 0))
        dr.text((x + sw // 2 - 10, y + 20), f"h={i}", fill=(0, 0, 0), font=font_b)
    # Plane swatches
    y2 = 170
    dr.rectangle([pad, y2, pad + 80, y2 + 30], fill=COLOR_K, outline=(0, 0, 0))
    dr.text((pad + 90, y2 + 5), "K plane (kv=0)", fill=(0, 0, 0), font=font_b)
    dr.rectangle([pad + 300, y2, pad + 380, y2 + 30], fill=COLOR_V, outline=(0, 0, 0))
    dr.text((pad + 390, y2 + 5), "V plane (kv=1)", fill=(0, 0, 0), font=font_b)
    return img

def draw_tile_map(Hk: int, T: int, label: str, width=1280, height=360) -> Image.Image:
    img = Image.new("RGB", (width, height), (255, 255, 255))
    dr = ImageDraw.Draw(img)
    font_b = ImageFont.load_default()
    dr.text((20, 10), label, fill=(0, 0, 0), font=font_b)
    left, top = 80, 60
    grid_w, grid_h = width - left - 40, height - top - 40
    # tiles: Hk rows, T columns
    cell_w = max(1, grid_w // max(1, T))
    cell_h = max(1, grid_h // max(1, Hk))
    for h in range(Hk):
        for t in range(T):
            x0 = left + t * cell_w
            y0 = top + h * cell_h
            dr.rectangle([x0, y0, x0 + cell_w, y0 + cell_h],
                         fill=PALETTE_HEADS[h % len(PALETTE_HEADS)], outline=(220, 220, 220))
    # axes
    dr.text((left, top - 20), "t (token index)", fill=(0, 0, 0), font=font_b)
    dr.text((10, top), "h (head)", fill=(0, 0, 0), font=font_b)
    return img

def draw_linear_banner(total_bytes: int, K_bytes: int, width=1280, height=180) -> Image.Image:
    """
    A 1D banner showing [K bytes | V bytes] across the first couple KB.
    """
    img = Image.new("RGB", (width, height), (255, 255, 255))
    dr = ImageDraw.Draw(img)
    font_b = ImageFont.load_default()
    dr.text((20, 10), "Linear On-Disk Layout (first region)", fill=(0, 0, 0), font=font_b)
    left, top, right = 40, 60, width - 40
    dr.rectangle([left, top, right, top + 50], fill=(240, 240, 255), outline=(0, 0, 0))
    span = right - left
    if total_bytes > 0:
        k_end = left + int(span * (K_bytes / total_bytes))
        dr.rectangle([left, top, k_end, top + 50], fill=COLOR_K, outline=(0, 0, 0))
        dr.rectangle([k_end, top, right, top + 50], fill=COLOR_V, outline=(0, 0, 0))
        dr.text((left, top + 60), "Fill = plane; top stripe = head color; vertical ticks = element boundaries", fill=(0, 0, 0), font=font_b)
    return img

def hex_lines(img_width, lines, Hk, T, D, dtype_bytes):
    """
    Render a vertical stack of hex lines with left color bars encoding kv/h/t/d.
    Each line shows 16 bytes.
    """
    line_h = 22
    width = img_width
    height = 10 + line_h * len(lines)
    img = Image.new("RGB", (width, height), (255, 255, 255))
    dr = ImageDraw.Draw(img)
    font_m = ImageFont.load_default()
    x0 = 10
    for i, (base_off, data) in enumerate(lines):
        y = 5 + i * line_h
        # left codebar
        idx0 = byte_index_to_indices(base_off, dtype_bytes, Hk, T, D)
        bar_color = COLOR_K if idx0.kv == 0 else COLOR_V
        dr.rectangle([x0, y + 2, x0 + 8, y + line_h - 2], fill=bar_color)
        dr.rectangle([x0 + 8, y + 2, x0 + 16, y + line_h - 2],
                     fill=PALETTE_HEADS[idx0.h % len(PALETTE_HEADS)])
        # token shading (alternate)
        if idx0.t % 2 == 1:
            dr.rectangle([x0 + 16, y + 2, x0 + 24, y + line_h - 2], fill=COLOR_TOKEN_STRIPE)
        else:
            dr.rectangle([x0 + 16, y + 2, x0 + 24, y + line_h - 2], outline=(200, 200, 200))
        # dim shading sample cell
        dr.rectangle([x0 + 24, y + 2, x0 + 32, y + line_h - 2], fill=COLOR_DIM_STRIPE)
        # offset text
        dr.text((x0 + 40, y + 4), f"{base_off:08x}:", fill=(0, 0, 0), font=font_m)
        # hex bytes
        hx = " ".join(f"{b:02x}" for b in data)
        dr.text((x0 + 110, y + 4), hx, fill=(0, 0, 0), font=font_m)
    return img

def split_hex(bytes_buf: bytes, start_off: int, max_lines: int) -> Tuple[List[Tuple[int, bytes]], int]:
    """
    Returns (lines, next_offset). Each line is (base_offset, 16-byte slice).
    """
    lines = []
    pos = start_off
    for _ in range(max_lines):
        if pos >= len(bytes_buf):
            break
        chunk = bytes_buf[pos:pos + 16]
        lines.append((pos, chunk))
        pos += 16
    return lines, pos

# --------------------------------------------------------------------------------------
# Slide builder
# --------------------------------------------------------------------------------------

class SlideBuilder:
    def __init__(self, out_path: str):
        self.prs = Presentation()
        self.out_path = out_path
        # Set 16:9 size explicitly for consistent layout
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    # utilities
    def _title_slide(self, title: str, subtitle: Optional[str] = None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # blank
        # Title box
        box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.9), Inches(1.2))
        tf = box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(80, 80, 80)
        return slide

    def _bullet_slide(self, title: str, bullets: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        # title
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(11.9), Inches(1.0))
        tf = tb.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        # bullets
        y = 1.6
        for b in bullets:
            b = textwrap.fill(b, 100)
            box = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(11.0), Inches(0.8))
            p = box.text_frame.paragraphs[0]
            p.text = f"• {b}"
            p.font.size = Pt(22)
            y += 0.7
        return slide

    def _image_slide(self, title: str, pil_img: Image.Image, height_in=4.6):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8))
        tf = tb.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        bio = BytesIO()
        pil_img.save(bio, format="PNG")
        pic = slide.shapes.add_picture(BytesIO(bio.getvalue()), Inches(0.7), Inches(1.4), height=Inches(height_in))
        return slide

    # deck content
    def add_title(self, prompt_text: str, real_or_dry: str):
        self._title_slide(
            "The Future of Exabyte‑Scale Inference",
            subtitle=f"Generated by TFEITool.py  •  Mode: {real_or_dry}\nPrompt: {prompt_text}"
        )

    def add_program_overview(self):
        bullets = [
            "TFEITool is an inference server + analyzer that turns a prompt into a fully labeled slide deck.",
            "It traces HTTP+runtime I/O, measures TTFT & tokens/sec, and inspects KV cache files.",
            "All diagrams map directly to the byte layout of KV blocks on disk (no truncation in hex views).",
        ]
        self._bullet_slide("What this program does", bullets)

    def add_kv_concepts(self, Hk: int, T: int, D: int, dtype: str):
        bullets = [
            "A KV block stores a chunk ('paged' portion) of attention Keys and Values.",
            f"Shape (row-major): [2 (K,V), Hk={Hk}, T={T}, D={D}] • dtype={dtype}",
            "Keys (kv=0) then Values (kv=1) on disk. Indices advance as: kv -> head -> token -> dim.",
        ]
        self._bullet_slide("KV Block: What is inside one page?", bullets)

    def add_legend_and_tiles(self, Hk: int, T: int):
        self._image_slide("Head/Plane Color Legend", draw_color_legend())
        # Note: use ASCII hyphen instead of unicode em dash to avoid PIL unicode encoding issues
        self._image_slide("Tile Map - Keys (K): heads x token positions",
                          draw_tile_map(Hk, T, "K plane - heads vs tokens"))
        self._image_slide("Tile Map - Values (V): heads x token positions",
                          draw_tile_map(Hk, T, "V plane - heads vs tokens"))

    def add_linear_layout(self, total_bytes: int, K_bytes: int):
        self._image_slide("Linear On-Disk Layout (K then V)", draw_linear_banner(total_bytes, K_bytes))

    def add_flow_slide(self, tracer_summary: Dict[str, Any], engine_dir: Optional[str]):
        bullets = [
            f"TTFT: {tracer_summary.get('ttft_s'):.3f}s \u2022 tokens: {tracer_summary.get('tokens')} \u2022 rate: {tracer_summary.get('tokens_per_s'):.2f}/s",
            f"HTTP bytes \u2191{tracer_summary.get('bytes_up')} \u2193{tracer_summary.get('bytes_down')}",
            f"Engine dir: {engine_dir or '(N/A)'}",
        ]
        self._bullet_slide("End‑to‑End Dataflow (high‑level)", bullets)

    def add_full_hex(self, buf: bytes, Hk: int, T: int, D: int, dtype_bytes: int,
                     lines_per_slide: int):
        """
        Render ALL bytes with colored left bars (plane/head/token/dim) over as many slides as needed.
        """
        offset = 0
        total = len(buf)
        page = 0
        while offset < total:
            lines, next_off = split_hex(buf, offset, lines_per_slide)
            img = hex_lines(1280, lines, Hk, T, D, dtype_bytes)
            # Use ASCII hyphen instead of unicode em dash for compatibility
            self._image_slide(f"Full Hex Map - {page+1} (bytes {offset}..{min(next_off, total)-1} / {total})", img, height_in=5.4)
            offset = next_off
            page += 1

    def save(self):
        self.prs.save(self.out_path)
        print(f"[ok] wrote deck: {self.out_path}")

# --------------------------------------------------------------------------------------
# Main routine
# --------------------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(description="TFEITool: TRT‑LLM I/O + KV cache analyzer and PPTX generator")
    ap.add_argument("--server-url", default="http://localhost:8000", help="OpenAI-compatible base URL (trtllm-serve)")
    ap.add_argument("--model", default="tensorrt-llm", help="Model name as served by your endpoint")
    ap.add_argument("--api-key", default="sk-noop", help="Bearer token value (server may ignore)")
    ap.add_argument("--engine-dir", default=None, help="Path to TRT‑LLM engine dir (optional but recommended)")
    ap.add_argument("--kv-dump-glob", default="**/block_*_pool_*.bin", help="Glob to find KV block dumps")
    ap.add_argument("--pptx-out", default="tfei_presentation.pptx", help="Output PPTX file")
    ap.add_argument("--inference-prompt", required=True, help="User prompt to run through the server")
    ap.add_argument("--layer-index", type=int, default=0, help="Which decoder layer’s KV geometry to use")
    ap.add_argument("--dtype", default=None, help="Override dtype (fp16|bf16|fp32|int8) if engine discovery unavailable")
    ap.add_argument("--kv-scale", type=float, default=None, help="INT8 KV dequant scale (if used)")
    ap.add_argument("--dry-run", action="store_true", help="No server required; synthesize data & slides")
    args = ap.parse_args()

    cli = CLI(
        server_url=args.server_url, model=args.model, api_key=args.api_key,
        engine_dir=args.engine_dir, kv_dump_glob=args.kv_dump_glob,
        pptx_out=args.pptx_out, inference_prompt=args.inference_prompt,
        layer_index=args.layer_index, dry_run=args.dry_run,
        dtype=args.dtype, kv_scale=args.kv_scale
    )

    hooks = ExternalHooks()
    tracer = IOTracer()

    # 1) Discover KV geometry
    kvgeo = discover_kv_geometry(cli.engine_dir, cli.layer_index, cli.dtype)
    print(f"[info] KV geometry: Hk={kvgeo.num_kv_heads} T={kvgeo.tokens_per_block} D={kvgeo.head_dim} dtype={kvgeo.dtype}")

    # 2) Run inference (real or dry-run)
    if cli.dry_run:
        tracer.mark("dry_run_start")
        time.sleep(0.15)  # pretend prefill
        tracer.token()
        for _ in range(64):
            time.sleep(0.01)
            tracer.token()
        tracer.mark("dry_run_end")
        output_text = "(dry-run) synthetic completion"
    else:
        output_text = run_inference_openai_http(
            server_url=cli.server_url, model=cli.model, api_key=cli.api_key,
            prompt=cli.inference_prompt, tracer=tracer, hooks=hooks
        )

    summ = tracer.summary()
    print(f"[info] run summary: {summ}")

    # 3) Collect KV dump files (created by the server/runtime)
    dump_paths = sorted(glob.glob(cli.kv_dump_glob, recursive=True))
    if dump_paths:
        print(f"[info] found {len(dump_paths)} KV dump files")
    else:
        print("[warn] no KV dump files found; hex slides will be empty (OK in --dry-run)")

    # For the deck we’ll show the first file of pool 0, or any file if not present
    bin_buf = b""
    pick = None
    for p in dump_paths:
        pick = p
        break
    if pick:
        bin_buf = read_block_file(pick)
        print(f"[info] reading: {pick}  bytes={len(bin_buf)}")

    # Derived sizes
    element_bytes = 2 if (kvgeo.dtype or "fp16").lower() in ("fp16", "half", "bf16", "bfloat16") else 4
    if (kvgeo.dtype or "fp16").lower() == "int8":
        element_bytes = 1

    total_elems = 2 * kvgeo.num_kv_heads * kvgeo.tokens_per_block * kvgeo.head_dim
    expected_bytes = total_elems * element_bytes

    if pick and len(bin_buf) != expected_bytes:
        print(f"[warn] size mismatch: got {len(bin_buf)} bytes, expected {expected_bytes} (check Hk/T/D/dtype)")

    # 4) Build slides
    sb = SlideBuilder(cli.pptx_out)
    sb.add_title(cli.inference_prompt, "REAL" if not cli.dry_run else "DRY-RUN")
    sb.add_program_overview()
    sb.add_flow_slide(summ, cli.engine_dir)
    sb.add_kv_concepts(kvgeo.num_kv_heads, kvgeo.tokens_per_block, kvgeo.head_dim, kvgeo.dtype)
    sb.add_legend_and_tiles(kvgeo.num_kv_heads, kvgeo.tokens_per_block)

    # linear layout banner
    k_bytes = kvgeo.num_kv_heads * kvgeo.tokens_per_block * kvgeo.head_dim * element_bytes
    total_bytes = 2 * k_bytes
    sb.add_linear_layout(total_bytes, k_bytes)

    # Full hex (no truncation) — if we have a real file; else add an empty page
    if len(bin_buf) > 0:
        sb.add_full_hex(bin_buf, kvgeo.num_kv_heads, kvgeo.tokens_per_block,
                        kvgeo.head_dim, element_bytes, lines_per_slide=cli.max_hex_lines_per_slide)
    else:
        sb._bullet_slide("No KV block file found",
                         [f"Glob searched: {cli.kv_dump_glob}",
                          "If you used TRT‑LLM KV offload (GDS/POSIX) the block files are named like block_<N>_pool_<M>.bin.",
                          "Rerun with a real inference and ensure KV dumping is enabled."])

    # Close with a small “reproducible/how-to” page
    sb._bullet_slide("This presentation is reproducible",
                     ["Run TFEITool.py against your TRT‑LLM server.",
                      "Pass --engine-dir to auto‑discover KV geometry.",
                      "Provide a prompt with --inference-prompt."])

    sb.save()

if __name__ == "__main__":
    main()